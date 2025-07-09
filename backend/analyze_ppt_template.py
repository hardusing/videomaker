#!/usr/bin/env python3
"""
分析PPT模板结构的脚本
"""
from pptx import Presentation
import os

def analyze_ppt_template():
    ppt_file = 'pptx/ストアドプロシージャ基礎.pptx'
    
    if not os.path.exists(ppt_file):
        print(f"❌ PPT文件不存在: {ppt_file}")
        return
    
    try:
        prs = Presentation(ppt_file)
        print(f"✅ PPT文件: {ppt_file}")
        print(f"📊 总页数: {len(prs.slides)}")
        print("=" * 60)
        
        for i, slide in enumerate(prs.slides, 1):
            print(f"\n📄 第{i}页:")
            print(f"   布局: {slide.slide_layout.name}")
            
            # 分析标题
            if slide.shapes.title:
                title_text = slide.shapes.title.text
                print(f"   标题: {title_text}")
                print(f"   标题位置: left={slide.shapes.title.left}, top={slide.shapes.title.top}")
                print(f"   标题大小: width={slide.shapes.title.width}, height={slide.shapes.title.height}")
            
            # 分析所有形状
            print(f"   形状数量: {len(slide.shapes)}")
            for j, shape in enumerate(slide.shapes):
                shape_type = shape.shape_type
                print(f"     形状{j+1}: {shape_type}")
                
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    text_content = shape.text_frame.text
                    if text_content.strip():
                        preview = text_content[:50].replace('\n', '\\n')
                        print(f"       文本: {preview}...")
                        print(f"       位置: left={shape.left}, top={shape.top}")
                        print(f"       大小: width={shape.width}, height={shape.height}")
                
                # 安全检查占位符
                try:
                    if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                        print(f"       占位符类型: {shape.placeholder_format.type}")
                except:
                    pass  # 忽略占位符检查错误
            
            print("-" * 40)
            
            # 只分析前3页作为示例
            if i >= 3:
                print("... (省略其余页面分析)")
                break
    
    except Exception as e:
        print(f"❌ 分析PPT时出错: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    analyze_ppt_template() 