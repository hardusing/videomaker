# videomaker/backend/app/utils/ppt_generator.py
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

def create_ppt_from_template(template_path: str, output_path: str, slides_data: list):
    """
    根据PPTX模板创建新的演示文稿，并应用自定义格式。
    会先清空模板中所有已存在的幻灯片，然后根据数据重新生成。

    :param template_path: 模板文件的路径
    :param output_path: 生成的PPTX文件的保存路径
    :param slides_data: 一个包含幻灯片数据的列表，每个元素是一个字典，
                        例如: [{'title': '标题', 'content': '内容'}, ...]
    """
    prs = Presentation(template_path)
    # 定义颜色
    title_color = RGBColor(77, 113, 179) # 精确颜色 #4D71B3
    black_color = RGBColor(0, 0, 0) # 黑色

    # --- 清空模板中所有已存在的幻灯片 ---
    sldIdLst = prs.slides._sldIdLst
    for i in range(len(sldIdLst) - 1, -1, -1):
        sldId = sldIdLst[i]
        prs.part.drop_rel(sldId.rId)
        del sldIdLst[i]
    # ------------------------------------

    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    for i, slide_data in enumerate(slides_data):
        if i == 0 and 'title' in slide_data:
            slide = prs.slides.add_slide(title_slide_layout)

            if len(slide.placeholders) > 1:
                try:
                    subtitle_shape = slide.placeholders[1]
                    sp_element = subtitle_shape._element
                    sp_element.getparent().remove(sp_element)
                except KeyError:
                    pass

            title_shape = slide.shapes.title
            title_shape.left = 0
            title_shape.top = 0
            title_shape.width = prs.slide_width
            title_shape.height = prs.slide_height

            title_shape.text = slide_data.get('title', '默认标题')
            text_frame = title_shape.text_frame
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            p_title = text_frame.paragraphs[0]
            p_title.alignment = PP_ALIGN.CENTER
            p_title.font.name = 'Microsoft YaHei'
            p_title.font.bold = True
            p_title.font.size = Pt(40)
            p_title.font.color.rgb = title_color

        else:
            slide = prs.slides.add_slide(content_slide_layout)
            
            title_shape = slide.shapes.title
            title_shape.left = Cm(2.01)
            title_shape.top = Cm(2.2)
            title_shape.width = prs.slide_width - (Cm(2.01) * 2)
            title_shape.height = Cm(2.19)

            p_title = title_shape.text_frame.paragraphs[0]
            p_title.text = slide_data.get('title', f'第 {i+1} 页')
            p_title.alignment = PP_ALIGN.LEFT
            p_title.font.name = 'Microsoft YaHei'
            p_title.font.bold = True
            p_title.font.size = Pt(20)
            p_title.font.color.rgb = title_color
            
            if len(slide.placeholders) > 1:
                body_shape = slide.placeholders[1]
                content = slide_data.get('content', '')

                if body_shape is not None:
                    try:
                        sp_element = body_shape._element
                        sp_element.getparent().remove(sp_element)
                    except Exception:
                        pass
                
                content_area_left = Cm(2.01)
                content_area_top = Cm(2.2 + 2.19 + 0.5)
                content_area_width = prs.slide_width - (Cm(2.01) * 2)
                content_area_height = prs.slide_height - content_area_top - Cm(2.01)

                content_blocks = re.split(r'(```.*?```)', content, flags=re.DOTALL)
                
                current_y = content_area_top
                for block_text in content_blocks:
                    if not block_text.strip(): continue
                    
                    is_code = block_text.startswith('```')
                    clean_text = block_text[3:-3].strip() if is_code else block_text.strip()
                    
                    box = slide.shapes.add_textbox(content_area_left, current_y, content_area_width, Cm(5)) # Start with a default height
                    tf = box.text_frame
                    tf.clear()
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = clean_text
                    
                    if is_code:
                        font = p.font
                        font.name = 'Source Code Pro'
                        font.size = Pt(12)
                        font.color.rgb = black_color
                        p.line_spacing = 1.0
                        fill = box.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(240, 240, 240)
                    else:
                        font = p.font
                        font.name = '游ゴシック'
                        font.size = Pt(16)
                        font.color.rgb = black_color
                        p.line_spacing = 1.5
                    
                    p.alignment = PP_ALIGN.LEFT
                    tf.auto_size = True
                    current_y += tf.height + Pt(10)


    output_dir = os.path.dirname(output_path)
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        
    prs.save(output_path)
    print(f"演示文稿已保存至: {output_path}")

if __name__ == '__main__':
    current_dir = os.path.dirname(os.path.abspath(__file__))
    backend_dir = os.path.join(current_dir, '..', '..')
    
    template_file = os.path.join(backend_dir, 'pptx', 'ppt模板.pptx')
    output_dir_path = os.path.join(backend_dir, 'outputs')
    
    slides_content = [
        {
            'title': 'ボックスモデルとレイアウトの基礎',
            'content': 'CSSレイアウトの核心を理解する'
        },
        {
            'title': '1. ボックスモデルの概要',
            'content': 'すべてのHTML要素は、ブラウザによって四角形の「ボックス」として描画されます。このボックスは4つの層から構成されており、CSSでレイアウトを調整する際の基本的な単位となります。\n\n- **内容 (Content)**: テキストや画像が表示される中心領域。\n- **内側余白 (Padding)**: 内容と枠線の間の透明なスペース。\n- **枠線 (Border)**: Paddingを囲む線。太さ、色、スタイルを指定可能。\n- **外側余白 (Margin)**: 枠線の外側にある透明なスペース。他の要素との距離を定義します。'
        },
        {
            'title': '2. Content・Padding・Border・Marginの違い',
            'content': 'CSSプロパティを使って、ボックスモデルの各層を具体的に制御します。これにより、要素の内部的な余白や、他の要素との間隔を自由に設定できます。\n\n以下のコードは、ボックスに内側余白、枠線、外側余白を同時に指定した例です。\n```.box {\n  padding: 10px; /* 内側余白を上下左右に10px */\n  border: 2px solid black; /* 2pxの黒い実線の枠 */\n  margin: 20px; /* 外側余白を上下左右に20px */\n}```'
        },
        {
            'title': '3. widthとheightの指定',
            'content': '`width` と `height` プロパティは、要素の「内容(Content)」領域のサイズをデフォルトで指定します。PaddingやBorderのサイズはこれに含まれないため、全体の大きさは `width` + `padding` + `border` となります。\n\nこの挙動は `box-sizing: border-box;` を指定することで変更でき、`width` と `height` が枠線まで含んだ全体のサイズを示すようになり、直感的なレイアウト作成が可能になります。\n```.box {\n  width: 300px; /* 内容領域の幅が300px */\n  height: 150px; /* 内容領域の高さが150px */\n}```'
        },
        {
            'title': '4. displayプロパティの基礎',
            'content': '`display` プロパティは、要素がどのように表示され、他の要素とどう影響しあうかを決定します。主な値は以下の通りです。\n\n- **block**: 幅が親要素いっぱいに広がり、前後に改行が入ります。(例: `<div>`, `<p>`)\n- **inline**: 必要な幅しか持たず、前後に改行が入りません。`width` や `height` は指定できません。(例: `<span>`, `<a>`)\n- **inline-block**: `inline` のように他の要素と並びますが、`block` のように `width` や `height` を指定できます。'
        },
        {
            'title': '5. positionプロパティの基本',
            'content': '`position` プロパティは、要素の配置方法を標準のドキュメントフローから変更する際に使用します。\n\n- **static**: デフォルト値。通常の流れに従って配置されます。\n- **relative**: 通常の位置を基準に、`top`, `right`, `bottom`, `left` で相対的に位置を調整できます。\n- **absolute**: 親要素の中で `position` が `static` 以外に指定された最も近いものを基準に絶対位置で配置されます。\n- **fixed**: ブラウザの表示領域（ビューポート）を基準に固定位置で配置され、スクロールしても動きません。'
        },
        {
            'title': '6. Flexboxの導入',
            'content': 'Flexboxは、1次元（縦一列または横一列）のレイアウトを簡単かつ柔軟に構築するための仕組みです。要素間のスペースの分配や、中央揃えなどを直感的に行うことができます。\n\nコンテナ要素に `display: flex;` を指定することで、その子要素（アイテム）がFlexboxのレイアウトに従うようになります。\n```.container {\n  display: flex;\n  /* アイテムを水平方向の中央に揃える */\n  justify-content: center;\n}```'
        },
        {
            'title': '7. Flexboxの主なプロパティ',
            'content': 'Flexboxはコンテナとアイテムにそれぞれプロパティを指定することで、多彩なレイアウトを実現します。\n\n**コンテナ側の主なプロパティ:**\n- `justify-content`: 主軸（横方向）の揃え方を指定 (例: `flex-start`, `center`, `space-between`)\n- `align-items`: 交差軸（縦方向）の揃え方を指定 (例: `flex-start`, `center`, `stretch`)\n- `flex-direction`: 主軸の方向を指定 (`row` または `column`)'
        },
        {
            'title': '8. Gridレイアウトの概要',
            'content': 'Gridレイアウトは、行と列から成る2次元の格子状レイアウトを構築するための強力なシステムです。Flexboxが1次元のレイアウトを得意とするのに対し、Gridはより複雑で大規模なレイアウト設計に適しています。\n\nコンテナに `display: grid;` を指定し、`grid-template-columns` や `grid-template-rows` で格子構造を定義します。\n```.grid-container {\n  display: grid;\n  /* 2つの列を作成。1列目は利用可能スペースの1/3、2列目は2/3を占める */\n  grid-template-columns: 1fr 2fr;\n}```'
        },
        {
            'title': '9. floatとclearの使い方',
            'content': '`float` は、主に画像などを横に配置し、後続のテキストをその周りに回り込ませるために使われる古典的なプロパティです。しかし、`float` を使うと親要素の高さが認識されなくなるなどのレイアウト崩れが起きやすいという欠点があります。\n\n`clear` プロパティは、`float` による回り込みを解除するために使用されます。\n```img {\n  float: left; /* 画像を左に寄せる */\n  margin-right: 10px;\n}\n\n.clearfix::after {\n  content: "";\n  display: block;\n  clear: both; /* floatを解除するおまじない */\n}```'
        },
        {
            'title': '10. レスポンシブレイアウトの考え方',
            'content': 'レスポンシブレイアウトとは、PC、タブレット、スマートフォンなど、異なる画面サイズのデバイスでウェブサイトが最適に表示されるように設計する考え方です。\n\nCSSの「メディアクエリ (`@media`)」を使用することで、特定の画面幅になったときに適用するスタイルを切り替えることができます。\n```/* 画面幅が768px以下の場合に適用 */\n@media (max-width: 768px) {\n  .container {\n    /* 例えば、横並びを縦並びに変更 */\n    flex-direction: column;\n  }\n}```'
        }
    ]
    
    first_page_title = slides_content[0].get('title', 'generated_presentation')
    safe_filename = "".join([c for c in first_page_title if c not in r'\\/:*?"<>|']) + ".pptx"
    output_file_path = os.path.join(output_dir_path, safe_filename)
    
    if not os.path.exists(template_file):
        print(f"错误：模板文件未找到，路径: {template_file}")
    else:
        create_ppt_from_template(template_file, output_file_path, slides_content) 