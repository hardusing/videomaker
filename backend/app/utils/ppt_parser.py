# utils/ppt_parser.py
import os
from pptx import Presentation

def extract_notes(pptx_path: str, output_dir: str) -> list:
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    prs = Presentation(pptx_path)
    page_texts = []

    for index, slide in enumerate(prs.slides, start=1):
        notes = slide.notes_slide.notes_text_frame.text if slide.notes_slide else ""
        output_file = os.path.join(output_dir, f"{index}.txt")
        with open(output_file, "w", encoding="utf-8") as f:
            f.write(notes)
        page_texts.append({"page": index, "content": notes})

    return page_texts
