#!/usr/bin/env python3
"""
基于ストアドプロシージャ基礎.pptx模板的内容替换脚本
支持13页PPT的内容替换，包括标题、文本、代码块等
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import re
import copy

class PPTContentReplacer:
    def __init__(self, template_path='pptx/ストアドプロシージャ基礎.pptx'):
        self.template_path = template_path
        self.prs = None
        
        # 定义样式配置
        self.styles = {
            'title_color': RGBColor(77, 113, 179),  # #4D71B3
            'text_color': RGBColor(0, 0, 0),        # 黑色
            'code_bg_color': RGBColor(240, 240, 240),  # 浅灰色背景
            'code_text_color': RGBColor(0, 0, 139),    # 深蓝色
            'title_font': 'Microsoft YaHei',
            'text_font': '游ゴシック', 
            'code_font': 'Consolas'
        }
    
    def load_template(self):
        """加载PPT模板"""
        if not os.path.exists(self.template_path):
            raise FileNotFoundError(f"模板文件不存在: {self.template_path}")
        
        self.prs = Presentation(self.template_path)
        print(f"✅ 已加载模板: {self.template_path}")
        print(f"📊 模板页数: {len(self.prs.slides)}")
        return self.prs
    
    def replace_content(self, slides_content):
        """
        替换PPT内容
        slides_content: 包含13页内容的列表，每个元素是字典
        格式: [
            {'title': '标题', 'content': '内容'},
            ...
        ]
        """
        if not self.prs:
            self.load_template()
        
        if len(slides_content) != 13:
            raise ValueError(f"内容必须包含13页，当前提供了{len(slides_content)}页")
        
        for i, slide_data in enumerate(slides_content):
            if i < len(self.prs.slides):
                self._replace_slide_content(i, slide_data)
            else:
                print(f"⚠️ 警告: 第{i+1}页超出模板范围")
    
    def _replace_slide_content(self, slide_index, slide_data):
        """替换单页内容"""
        slide = self.prs.slides[slide_index]
        title = slide_data.get('title', '')
        content = slide_data.get('content', '')
        
        print(f"🔄 处理第{slide_index+1}页: {title}")
        
        # 替换标题
        if slide.shapes.title and title:
            slide.shapes.title.text = title
            self._apply_title_style(slide.shapes.title, slide_index == 0)
        
        # 处理内容
        if content:
            self._replace_content_area(slide, content, slide_index)
    
    def _apply_title_style(self, title_shape, is_first_page=False):
        """应用标题样式"""
        if title_shape.text_frame:
            p = title_shape.text_frame.paragraphs[0]
            p.font.name = self.styles['title_font']
            p.font.bold = True
            p.font.color.rgb = self.styles['title_color']
            
            if is_first_page:
                # 首页标题样式
                p.font.size = Pt(40)
                p.alignment = PP_ALIGN.CENTER
                title_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            else:
                # 内容页标题样式
                p.font.size = Pt(20)
                p.alignment = PP_ALIGN.LEFT
    
    def _replace_content_area(self, slide, content, slide_index):
        """替换内容区域"""
        # 删除除标题外的所有文本框
        shapes_to_remove = []
        for shape in slide.shapes:
            if (shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or 
                shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER) and shape != slide.shapes.title:
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            try:
                sp_element = shape._element
                sp_element.getparent().remove(sp_element)
            except:
                pass
        
        # 添加新的内容
        if slide_index == 0:
            # 首页不需要额外内容
            return
        
        # 定义内容区域
        content_left = Cm(2.01)
        content_top = Cm(2.2 + 2.19 + 0.5)  # 标题下方
        content_width = self.prs.slide_width - (Cm(2.01) * 2)
        
        # 解析内容块（代码块和普通文本）
        content_blocks = self._parse_content_blocks(content)
        
        current_y = content_top
        for block in content_blocks:
            if block['type'] == 'code':
                current_y = self._add_code_block(slide, block['content'], content_left, current_y, content_width)
            else:
                current_y = self._add_text_block(slide, block['content'], content_left, current_y, content_width)
    
    def _parse_content_blocks(self, content):
        """解析内容块，分离代码和文本"""
        blocks = []
        
        # 使用正则表达式分离代码块和普通文本
        parts = re.split(r'(```[\s\S]*?```)', content)
        
        for part in parts:
            part = part.strip()
            if not part:
                continue
                
            if part.startswith('```') and part.endswith('```'):
                # 代码块
                code_content = part[3:-3].strip()
                blocks.append({'type': 'code', 'content': code_content})
            else:
                # 普通文本块
                blocks.append({'type': 'text', 'content': part})
        
        return blocks
    
    def _add_text_block(self, slide, text_content, left, top, width):
        """添加文本块"""
        # 创建文本框
        text_box = slide.shapes.add_textbox(left, top, width, Cm(1))
        text_frame = text_box.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        
        # 处理文本内容
        lines = text_content.split('\n')
        for i, line in enumerate(lines):
            line = line.strip()
            if not line:
                continue
                
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            # 处理不同类型的文本
            if line.startswith('**') and line.endswith('**'):
                # 粗体文本
                p.text = line[2:-2]
                p.font.bold = True
                p.font.size = Pt(16)
            elif line.startswith('- '):
                # 列表项
                p.text = line.replace('- ', '• ')
                p.font.size = Pt(16)
                p.level = 0
            elif line.startswith('  - '):
                # 二级列表项
                p.text = line.replace('  - ', '  • ')
                p.font.size = Pt(14)
                p.level = 1
            else:
                # 普通文本
                p.text = line
                p.font.size = Pt(16)
            
            p.font.name = self.styles['text_font']
            p.font.color.rgb = self.styles['text_color']
            p.alignment = PP_ALIGN.LEFT
        
        # 手动调整高度而不是使用auto_size
        # text_frame.auto_size = True  # 这行会导致错误
        return top + text_box.height + Pt(10)
    
    def _add_code_block(self, slide, code_content, left, top, width):
        """添加代码块"""
        # 创建代码文本框
        code_box = slide.shapes.add_textbox(left, top, width, Cm(1))
        
        # 设置背景色
        fill = code_box.fill
        fill.solid()
        fill.fore_color.rgb = self.styles['code_bg_color']
        
        # 设置文本
        text_frame = code_box.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.margin_left = Pt(10)
        text_frame.margin_right = Pt(10)
        text_frame.margin_top = Pt(8)
        text_frame.margin_bottom = Pt(8)
        
        p = text_frame.paragraphs[0]
        p.text = code_content
        p.font.name = self.styles['code_font']
        p.font.size = Pt(12)
        p.font.color.rgb = self.styles['code_text_color']
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.0
        
        # 手动调整高度而不是使用auto_size
        # text_frame.auto_size = True  # 这行会导致错误
        return top + code_box.height + Pt(15)
    
    def save(self, output_path):
        """保存PPT文件"""
        if not self.prs:
            raise ValueError("没有加载的PPT文件")
        
        # 确保输出目录存在
        output_dir = os.path.dirname(output_path)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir)
        
        self.prs.save(output_path)
        print(f"✅ PPT已保存: {output_path}")
        return output_path

def create_ppt_from_content(slides_content, output_path, template_path='pptx/ストアドプロシージャ基礎.pptx'):
    """
    便捷函数：从内容创建PPT
    
    参数:
    slides_content: 13页内容的列表
    output_path: 输出文件路径
    template_path: 模板文件路径
    """
    replacer = PPTContentReplacer(template_path)
    replacer.load_template()
    replacer.replace_content(slides_content)
    return replacer.save(output_path)

# 示例用法
if __name__ == "__main__":
    # 示例内容
    example_content = [
        {'title': '新课程标题'},
        {'title': '第一章 概述', 'content': '''这是第一章的内容。

**主要要点：**
- 要点1
- 要点2
- 要点3

```SELECT * FROM table_name;```

这是代码块后的说明文字。'''},
        {'title': '第二章 基础知识', 'content': '第二章的内容...'},
        {'title': '第三章', 'content': '第三章的内容...'},
        {'title': '第四章', 'content': '第四章的内容...'},
        {'title': '第五章', 'content': '第五章的内容...'},
        {'title': '第六章', 'content': '第六章的内容...'},
        {'title': '第七章', 'content': '第七章的内容...'},
        {'title': '第八章', 'content': '第八章的内容...'},
        {'title': '第九章', 'content': '第九章的内容...'},
        {'title': '第十章', 'content': '第十章的内容...'},
        {'title': '第十一章', 'content': '第十一章的内容...'},
        {'title': '总结', 'content': '课程总结内容...'}
    ]
    
    try:
        output_file = create_ppt_from_content(
            example_content, 
            'pptx/新课程_示例.pptx'
        )
        print(f"🎉 示例PPT创建完成: {output_file}")
    except Exception as e:
        print(f"❌ 创建PPT时出错: {e}")
        import traceback
        traceback.print_exc() 